from pptx import Presentation
import os
import get_bibles
import copy
from pptx.util import Pt
from pptx.dml.color import RGBColor

def read_pptx(pptx_file):
    """
    读取现有的PPTX文件并返回所有内容信息
    
    Args:
        pptx_file: PPTX文件路径
    
    Returns:
        dict: 包含PPT所有信息的字典
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return None
    
    prs = Presentation(pptx_file)
    ppt_info = {
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    # 遍历所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': slide_num,
            'title': '',
            'shapes': []
        }
        
        # 获取标题
        if slide.shapes.title:
            slide_info['title'] = slide.shapes.title.text
        
        # 遍历所有形状
        for shape_num, shape in enumerate(slide.shapes):
            shape_info = {
                'shape_number': shape_num,
                'type': str(shape.shape_type),
                'has_text': shape.has_text_frame,
                'text': ''
            }
            
            # 获取文本内容
            if shape.has_text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_parts.append(run.text)
                shape_info['text'] = ''.join(text_parts)
            
            slide_info['shapes'].append(shape_info)
        
        ppt_info['slides'].append(slide_info)
    
    return ppt_info


def update_pptx_text(pptx_file, output_file, replacements):
    """
    修改PPTX文件中的文字
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        replacements: 字典，格式 {'旧文字': '新文字'}
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 遍历所有幻灯片
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 进行替换
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
    
    # 保存
    prs.save(output_file)
    print(f"PPT文件已保存: {output_file}")
    return True


def print_pptx_info(ppt_info):
    """
    打印PPT信息
    
    Args:
        ppt_info: read_pptx函数返回的信息字典
    """
    if not ppt_info:
        return
    
    print(f"总共 {ppt_info['slide_count']} 页")
    print("=" * 60)
    
    for slide_info in ppt_info['slides']:
        print(f"\n第 {slide_info['slide_number']} 页")
        print(f"标题: {slide_info['title']}")
        print(f"形状数量: {len(slide_info['shapes'])}")       
        for shape_info in slide_info['shapes']:
            
            if shape_info['has_text'] and shape_info['text']:
                print(f"  - 文本: {shape_info['text']}...")


def print_pptx_page(ppt_info, page_number):
    """
    打印PPT信息
    
    Args:
        ppt_info: read_pptx函数返回的信息字典
    """
    if not ppt_info:
        return
    
    print(f"总共 {ppt_info['slide_count']} 页")
    print("=" * 60)
    
    for slide_info in ppt_info['slides']:
        if slide_info['slide_number'] != page_number:
            continue
        print(f"\n第 {slide_info['slide_number']} 页")
        print(f"标题: {slide_info['title']}")
        print(f"形状数量: {len(slide_info['shapes'])}")       
        for shape_info in slide_info['shapes']:
            
            if shape_info['has_text'] and shape_info['text']:
                print(f"  - 文本: {shape_info['text']}...")


def update_slide_text(pptx_file, output_file, slide_number, replacements):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        replacements: 字典，格式 {'旧文字': '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    # 遍历该页的所有形状
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # 进行替换
                    for old_text, new_text in replacements.items():
                        #print(f"Before {run.text} → {run.text.replace(old_text, new_text)}")
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                        #print(f"After {run.text}")
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True


def update_multiple_slides(pptx_file, output_file, slide_replacements):
    """
    批量修改多页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_replacements: 字典，格式 {页码: {'旧文字': '新文字'}}
        
    Example:
        slide_replacements = {
            1: {'标题': '新标题'},
            2: {'内容': '新内容'},
            3: {'2025': '2026'}
        }
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 遍历需要修改的页
    for slide_number, replacements in slide_replacements.items():
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 获取指定页
        slide = prs.slides[slide_number - 1]
        
        # 遍历该页的所有形状
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 进行替换
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
        
        print(f"已修改第 {slide_number} 页")
    
    # 保存
    prs.save(output_file)
    print(f"所有修改完成，文件已保存: {output_file}")
    return True


def delete_slide(pptx_file, output_file, slide_number):
    """
    删除指定页
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 要删除的页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取要删除的幻灯片
    rId = prs.slides._sldIdLst[slide_number - 1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_number - 1]
    
    # 保存
    prs.save(output_file)
    print(f"已删除第 {slide_number} 页，文件已保存: {output_file}")
    return True


def delete_slides(pptx_file, output_file, slide_numbers):
    """
    批量删除多页
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_numbers: 要删除的页码列表（从1开始），如 [2, 5, 7]
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 从大到小排序，从后往前删除，避免索引变化
    slide_numbers_sorted = sorted(slide_numbers, reverse=True)
    
    for slide_number in slide_numbers_sorted:
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 删除幻灯片
        rId = prs.slides._sldIdLst[slide_number - 1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_number - 1]
        print(f"已删除第 {slide_number} 页")
    
    # 保存
    prs.save(output_file)
    print(f"所有删除完成，文件已保存: {output_file}")
    return True


def duplicate_slides(pptx_file, output_file, slide_numbers):
    """
    批量复制多页并插入到各自后面
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_numbers: 要复制的页码列表（从1开始），如 [2, 5]
    
    Returns:
        bool: 是否成功
    """
    import copy
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 从大到小排序，从后往前处理，避免索引变化
    slide_numbers_sorted = sorted(slide_numbers, reverse=True)
    
    for slide_number in slide_numbers_sorted:
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 获取要复制的幻灯片
        source_slide = prs.slides[slide_number - 1]
        
        # 获取布局
        slide_layout = source_slide.slide_layout
        
        # 创建新幻灯片
        new_slide = prs.slides.add_slide(slide_layout)
        
        # 深度复制所有形状
        for shape in source_slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        # 移动到正确位置
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(slide_number, slides[-1])
        
        print(f"已在第 {slide_number} 页后插入副本")
    
    # 保存
    prs.save(output_file)
    print(f"所有复制完成，文件已保存: {output_file}")
    return True


def show_structure_one_page(pptx_file, slide_number):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        indexed_replacements: 字典，格式 {索引: '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]  
    
    for i in range(len(slide.shapes)):
        shape = slide.shapes[i]
        print(f"Shape index: {i}")
        if slide.shapes[i].has_text_frame:
            for j in range(len(slide.shapes[i].text_frame.paragraphs)):
                paragraph = slide.shapes[i].text_frame.paragraphs[j]
                print(f"  Paragraph index: {j}")
                for k in range(len(paragraph.runs)):
                    run = paragraph.runs[k]
                    print(f"     text index : {k} : {run.text}", end="|\n")


def duplicate_slide(pptx_file, output_file, slide_number):
    """
    复制指定页并插入到该页后面
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 要复制的页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    import copy
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取要复制的幻灯片
    source_slide = prs.slides[slide_number - 1]
    
    # 获取布局
    slide_layout = source_slide.slide_layout
    
    # 创建新幻灯片
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 深度复制所有形状
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # 移动到正确位置（紧跟在原页面后）
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(slide_number, slides[-1])
    
    # 保存
    prs.save(output_file)
    print(f"已在第 {slide_number} 页后插入副本，文件已保存: {output_file}")
    return True


def swap_slides(pptx_file, output_file, slide_num1, slide_num2):
    """
    交换两个幻灯片的位置
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_num1: 第一个页码（从1开始）
        slide_num2: 第二个页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_num1 < 1 or slide_num1 > len(prs.slides):
        print(f"错误：页码 {slide_num1} 超出范围（共 {len(prs.slides)} 页）")
        return False
    if slide_num2 < 1 or slide_num2 > len(prs.slides):
        print(f"错误：页码 {slide_num2} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    if slide_num1 == slide_num2:
        print("错误：两个页码不能相同")
        return False
    
    # 转换为0开始的索引
    idx1 = slide_num1 - 1
    idx2 = slide_num2 - 1
    
    # 获取XML中的slides列表
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    
    # 交换位置
    slides[idx1], slides[idx2] = slides[idx2], slides[idx1]
    
    # 清空并重新添加
    for slide in list(xml_slides):
        xml_slides.remove(slide)
    
    for slide in slides:
        xml_slides.append(slide)
    
    # 保存
    prs.save(output_file)
    print(f"已交换第 {slide_num1} 页和第 {slide_num2} 页，文件已保存: {output_file}")
    return True


def insert_fullscreen_video_slide(pptx_file, output_file, video_path, insert_position=None):
    """
    插入一个新的全屏视频幻灯片
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        video_path: 视频文件路径
        insert_position: 插入位置（从1开始），如果为None则在末尾添加
    
    Returns:
        bool: 是否成功
    """
    from pptx.util import Inches
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    if not os.path.exists(video_path):
        print(f"错误：找不到视频文件 {video_path}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 添加一个空白幻灯片（使用空白布局）
    blank_slide_layout = prs.slide_layouts[6]  # 6通常是空白布局
    new_slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加全屏视频
    # 视频位置：左上角(0,0)，尺寸：填满整个幻灯片
    left = Inches(0)
    top = Inches(0)
    width = slide_width
    height = slide_height
    
    # 插入视频
    movie = new_slide.shapes.add_movie(
        video_path,
        left, top, width, height,
        poster_frame_image=None,  # 不使用海报帧，使用视频第一帧
        mime_type='video/mp4'
    )
    
    # 如果指定了插入位置，则移动到该位置
    if insert_position is not None:
        if insert_position < 1 or insert_position > len(prs.slides):
            print(f"错误：插入位置 {insert_position} 超出范围（共 {len(prs.slides)} 页）")
            return False
        
        # 获取XML中的slides列表
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        
        # 移动新添加的幻灯片（最后一个）到指定位置
        xml_slides.remove(slides[-1])
        xml_slides.insert(insert_position - 1, slides[-1])
    
    # 保存
    prs.save(output_file)
    position_str = f"第 {insert_position} 页" if insert_position else "末尾"
    print(f"已在 {position_str} 插入全屏视频幻灯片，文件已保存: {output_file}")
    return True


def set_pptx_page_texts(pptx_file, output_file, slide_number, replacements):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        replacements: 字典，格式 {'旧文字': '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    # 遍历该页的所有形状
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for origin_text, change_text in replacements.items():
                            if origin_text in run.text:
                                print(f"{run.text} → {change_text}", end="\n")
                                run.text = run.text.replace(origin_text, change_text)
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True


def set_pptx_page_texts_by_slides_shapes_index(pptx_file, output_file, slide_number, replacements, size=False, resize=33, color=None):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        indexed_replacements: 字典，格式 {索引: '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    for shape_index, run_replacements in replacements.items():
        if not slide.shapes.__getitem__(shape_index).has_text_frame:
            print(f"错误：形状索引 {shape_index} 不包含文本框")
            return False
        shape = slide.shapes[shape_index]
        for paragraph_index in run_replacements.keys():
            paragraph = shape.text_frame.paragraphs[paragraph_index] 
            new_texts_index = run_replacements[paragraph_index]
            print(f"Shape {shape_index} paragraph {paragraph_index}")
            for run_index, new_text in new_texts_index.items():
                if run_index < len(paragraph.runs):
                    print(f" original text {paragraph.runs[run_index].text} new text: {new_text}")
                    paragraph.runs[run_index].text = new_text
                    paragraph.runs[run_index].font.bold = True
                    if size:
                        paragraph.runs[run_index].font.size = Pt(resize)
                    if color:
                        paragraph.runs[run_index].font.color.rgb = color
                else:
                    print(f" append new text on {run_index}: {new_text}")
                    new_run = paragraph.add_run()
                    new_run.text = " " + new_text
                    # 新增行时，字体加粗、字号30pt，字体固定为STXingkai
                    new_run.font.bold = True
                    new_run.font.size = Pt(20)
                    new_run.font.name = "STXingkai"
    
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True

def set_new_time(output_file, acceuil, date, page_to_modify=1):
    """
    page_to_modify = 1

    #date = "05/07/2026"
    
    """
    accueil = f"接待 {acceuil}"
    
    replacements = {0: {4: {0: "", 1: f"                  {date}", 2: " 15h-17h"}}, 1: {0: {8: f"。\n\n                                {accueil}"}}}

    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)
    return output_file

def set_daidao(page_to_modify, output_file, daidao_text, resize=21):
    """
    3.为教会的慕道友:金辉、张娟娟、丁建波、信祖生、白建亮、王子健、刘克辉，卜香峰，愿神的灵感动他们，让他们得着宝贵的救恩。\n4.为教会还在等待申请或延期居留的弟兄姐妹们祷告，求神预备，在这件事情上彰显神的荣耀。\n4.为那些还在找工作的弟兄姐妹们祷告。求神预备，并赐给他们合适的工作，能够安心生活在法国。\n5.为陈忠勇弟兄，宋立忠弟兄祷告，求神医治他们的身体，使他们快快的得到康复。\n6.徐霞姐妹家的小旋风belle 飞丢了一周 祈祷主把它引领回到家
    """
    replacements = {2: {0: {0: "教会代祷事项报告", 1: ""},
                        2: {0: "1", 1: daidao_text[0], 2: ""},
                        3: {0: "2", 1: daidao_text[1], 2: "", 3: "", 4: "", 5: "", 6: "", 7: "", 
                            8: "3", 9: daidao_text[2], 10: "", 11: "", 
                            12: "4", 13: daidao_text[3], 
                            14: "5", 15: daidao_text[4],
                            16: "6", 17: daidao_text[5], 18: "", 19: "", 
                            20: "7", 21: daidao_text[6], 
                            22: "8", 23: daidao_text[7], 
                            24: "9", 25: daidao_text[8], 
                            26: "10", 27: daidao_text[9]},
                    }}
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements, size=True, resize=resize, color=RGBColor(0, 0, 0))
    return output_file

def set_xuanzhao(output_file, name_linhui, titre, index, xuanzhao_text, page_to_modify=2):

    replacements = {0: {0: {2: titre, 3: "", 4: index, 5: ""}}, 1: {3: {1: name_linhui}, 1: {1: f"{xuanzhao_text}"}}}
    # update_slide_text(output_file, output_file, page_to_modify, {old_name: new_name})
    #set_pptx_page_texts(output_file, output_file, page_to_modify, replacements) 
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)
    return output_file

def set_jinbai_shiban(output_file, jinbai_text, shiban_title,shiban_text, if_shiban=True):
    page_to_modify = 3
    replacements = {1: {4: {1: f"{jinbai_text}"}}}
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)

    if if_shiban:
        page_to_modify = 7
        replacements = {0: {0: {1: "诗班献诗"}}, 1: {0: {0: f"                      {shiban_title}"}}, 2: {0: {0: f"{shiban_text}", 1: ""}}}
        set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)
    return output_file

def add_music(output_file, repository_music, page_huiyin, if_shiban=True, add_hui_ying=True):
    # musics 
    pages_music = [4,5,6]  # 假设音乐幻灯片是第4到第6页
    delete_slides(output_file, output_file, pages_music)  # 删除原有的音乐页
    
    for i in range(0, len(pages_music)):
        index = pages_music[i] - 3
        video_file = f"{repository_music}\\{index}.mp4"  # 修改为实际视频文件路径
        insert_fullscreen_video_slide(output_file, output_file, video_file, insert_position=(pages_music[i]))

    if if_shiban:
        pages_music = [8]  # 假设音乐幻灯片是第4到第6页
        delete_slides(output_file, output_file, pages_music)  # 删除原有的音乐页
        
        for i in range(0, len(pages_music)):
            index = pages_music[i] - 3
            video_file = f"{repository_music}\\{index}.mp4"  # 修改为实际视频文件路径
            insert_fullscreen_video_slide(output_file, output_file, video_file, insert_position=(pages_music[i]))

    if add_hui_ying:
        pages_music = [page_huiyin]  # 假设音乐幻灯片是第4到第6页
        delete_slides(output_file, output_file, pages_music)  # 删除原有的音乐页
        
        for i in range(0, len(pages_music)):
            video_file = f"{repository_music}\\3.mp4"  # 修改为实际视频文件路径
            insert_fullscreen_video_slide(output_file, output_file, video_file, insert_position=(pages_music[i]))
    
    return output_file


def zhirizhengdao(output_file, name_zhengdao, name_linhui, titre, index, page_to_modify=11):
    replacements = {1: {0: {0: "主日证道", 7: f"{titre}", 8: "", 9: "", 10: f"{index}", 11: "", 12: f"{name_zhengdao} 证道, {name_linhui} 回应"}}}
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)
    return output_file



if __name__ == "__main__":
    # 示例1：读取PPT信息
    filename = "template"

    template_repo = f'D:\\副业赚钱\\教会事务\\Template\\{filename}.pptx'
    repository = os.path.dirname(os.path.abspath(__file__))
    repository_music = os.path.join(os.path.dirname(repository), 'Template', 'musics')
    print(f"当前路径: {repository}")

    pptx_file = f"{repository}\\{filename}.pptx"
    output_file = f"{repository}\\{filename}.pptx"

    info = read_pptx(output_file)

    current_date = "23/08/2026"
    daidao_text = [
        "信望爱基督之家向法国政府申请协会和开银行的事都已经通过了，感谢主，请继续为在教会附近申请发放福音单张的事代祷，也愿神亲自成就，荣耀归主！",
        "周国莲姊妹从2026年6月1号起在信望爱之家实习传道代祷，求神赐智慧和能力，在她身上显明神的心意。\n",
"为9月份第一周受洗的弟兄姊妹，卜祥峰、张娟娟、丁建波、王锴、孙浩然、徐彦彬（待定）祷告，求神坚固他们的心。\n",
"为慕道的弟兄姐妹（金辉、王子健、白建亮、信祖生、刘克辉）祷告，愿神的灵感动他们，选择那上好的福份。\n",
 "为参与在爱心之家的同工valentina祷告，最近她在约旦服侍，求神与她同在，赐智慧和能力，有神迹伴随，见证神的荣耀！\n",
 "为在软弱中的肢体代祷，求神坚固他们的信心，重新回到神的里面，为陈忠勇、宋立忠两位弟兄祷告，求神医治他们的身体，使他们全然得康复;为卜祥峰弟兄祷告，求神医治的大能临到他的双膝，使他能行走不疼痛，荣耀归主！为吴兴隆弟兄祷告，求神医治他的左手\n",
 "为教会有需要找工作的弟兄姊妹们祷告，求神预备适合他们的工作，能够安心生活在法国。\n",
 "为教会在申请或延期、等待居留的弟兄姊妹（邸雪岩、卜祥峰、周国莲）祷告，在这件事情上看到神的荣耀。\n",
 "为在中东的宣教士们和他们的孩子来祷告，按着神荣耀的丰富，赐给他们一切所需用的都充足。\n",
 "为上法语课的老师和学生们祷告，通过学习，愿意更多的委身。\n",
 "为主日的讲台祷告，求神的灵大大充满，带下智慧和能力，荣耀归主！"]


    # 1 时间 + 接待
    #output_file = set_new_time(output_file, "巩象学弟兄", current_date, page_to_modify=1)
    
    # 代祷事项
    #output_file = set_daidao(12, output_file, daidao_text)

    
    # 领会
    page_to_modify = 2

    name_linhui="周国莲宣教士"
    titre = "以弗所书"
    index = "2：10"
    xuanzhao_text = "我们原是他的工作，在基督耶稣里造成的，为要叫我们行善，就是神所预备叫我们行的。"
    xuanzhao = (name_linhui, titre, index, xuanzhao_text)

    #show_structure_one_page(output_file, page_to_modify)
    # 宣召经文
    #output_file = set_xuanzhao(output_file, xuanzhao[0], xuanzhao[1], xuanzhao[2], xuanzhao[3], page_to_modify=2)

    # 敬拜
    jinbao_text = "徐霞姐妹 韩翠英姐妹"
    shiban_title = "炼我俞精"
    shiban_text = "徐霞姐妹, 巩象学弟兄"
    show_structure_one_page(output_file, 3)
    #output_file = set_jinbai_shiban(output_file, jinbao_text, shiban_title, shiban_text, if_shiban=False)

    # musics 
    #output_file = add_music(output_file, repository_music, if_shiban=False, page_huiyin=14, add_hui_ying=False)  # 假设音乐幻灯片是第13页


    # 主日证道
    name_zhengdao = "吳兴隆弟兄"
    titre = "谦卑与顺服"
    index = "腓立比书 2:1-30"
    page_zhuri = 12
    #output_file = zhirizhengdao(output_file, name_zhengdao, name_linhui, titre, index, page_to_modify=page_zhuri)


    # 圣餐
    delete_slides(output_file, output_file, [10])  # 删除原有的圣餐页
    #duplicate_slide(output_file, output_file, 14)  # 复制圣餐页作为模板
    page_to_modify = 15
    #show_structure_one_page(output_file, page_to_modify)

    replacements = {
        3: {0: {0: "", 1: "", 2: "", 3: "", 4: "", 5: "", 6: "", 7: "", 8: "", 9: "                            圣餐"}},
        4: {
            0: {0: "", 1: "", 2: ""},
            1: {0: "" , 2: ""},
            2: {0: "" , 2: ""},
            3: {0: "", 1: "耶稣说：“我就是生命的粮。到我这里来的，必定不饿；信我的，也必永远不渴。”", 2: ""},
            4: {0: "", 1: "", 2: ""},
            5: {0: "", 1: "", 2: ""},
            6: {0: "", 1: "", 2: ""},
            7: {0: "", 1: "", 2: ""},
            8: {0: "", 1: "", 2: ""}
        }
    }
    #set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements, size=True, resize=50, color=RGBColor(0, 0, 0))

    #insert_fullscreen_video_slide(output_file, output_file, f"{repository_music}\\6.mp4", insert_position=16)


    # 事工表
    page_to_modify = 18
    #show_structure_one_page(output_file, page_to_modify)
    #duplicate_slide(output_file, output_file, page_to_modify)
    replacements = {2: {2: {
                            4: "法语课 - 星期天", 
                            5: "13:30-14:30",
                            },
                        4: {
                            1: "旷野小组 – 星期三",
                            2: "20:30-22:30"
                        },
                        5: {
                            0: "                                                                                 ",
                            1: "荒漠小组 – 星期四",
                            2: "20:00-21:30 (线上)"
                        },
                        6: {
                            1: "祷告会 - 星期六",
                            2: "16:00-18:00",
                            3: "",
                            4: ""
                        },
                        7: {
                            0: "",
                            1: "",
                            2: "",
                            3: "",
                            4: ""
                        }
                        }
        }
    #set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements)

    #delete_slides(output_file, output_file, list(range(10, 13)))  # 删除多余的事工表页，保留第一页事工表页


    # comptable
    '''
四月份财务报表
上月余额 5177.16
本月奉献收入 1025

爱心扶助 100（陈弟兄）
教会厕所修理费 205.75
爱宴费用 165.85
宣教资助 300 （国莲）
爱心资助 300（徐霞）
法文教学费用 100
外来证道费用 70
本月余额 4961.06

5177.16 + 1025 - 4961.06 = 1241.1

    '''

    # 财务报表
    #duplicate_slide(output_file, output_file, 16)  # 复制第一页财务页作为模板
    page_to_modify = 16
    month = 7
    #show_structure_one_page(output_file, page_to_modify)
    #duplicate_slide(output_file, output_file, page_to_modify - 1)
    replacements = {3: {0: {4: "", 6: "", 8: "", 9: ""}}, 4: {1: {0: "", 1: "", 2: f"信望爱基督之家{month}月财务报告"}, 2: {0: f"{month - 1}月份余额 5247.02", 1: "", 2: ""}, 
                                                              3: {0: f"{month}收入", 1: "奉献收入 1085欧元", },
                                                              4: {0: f"{month}支出 1015.47 欧元", 1: "奉献总收入1085欧元\n爱心助教费用600欧元 爱心扶助 100欧元\n法文班教师 100欧元\n打印机墨水 74.73欧元\n洗礼大毛巾 29.24欧元\n爱宴总费用 111.5欧元", 2: ""},
                                                              5: {0: f"{month}月份余额", 1: " 5316.55欧元", 2: ""}}}
    #set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements, size=True, resize=33, color=RGBColor(0, 0, 0))

    #delete_slides(output_file, output_file, list(range(16, 20))+[21, 22])  # 删除多余的财务页，保留第一页财务页
  
    # ========== 经文页面（新方法）==========
    # 使用新函数设置经文页面
    
    # 第1页经文：路加福音 8:1-5（5行）
    
    #delete_slides(output_file, output_file, [7,7,7, 12, 12, 13])  # 删除多余的经文页，保留第一页经文页
    #delete_slides(output_file, output_file, [11 for i in range(0, 1)])
    #duplicate_slide(output_file, output_file, 10)  # 复制第一页经文页作为模板
    page_to_modify = 10
    title = "腓立比书" 
    chapter = 2
    index_text = get_bibles.indexes[title]

    #show_structure_one_page(output_file, page_to_modify)


    texts = [
            [title, chapter, 1, 5, get_bibles.get_bible_verses(index_text, chapter, 1, 5)],
            [title, chapter, 6, 10, get_bibles.get_bible_verses(index_text, chapter, 6, 10)],
            [title, chapter, 11, 15, get_bibles.get_bible_verses(index_text, chapter, 11, 15)],
            [title, chapter, 16, 20, get_bibles.get_bible_verses(index_text, chapter, 16, 20)],
            [title, chapter, 21, 25, get_bibles.get_bible_verses(index_text, chapter, 21, 25)],
            [title, chapter, 26, 30, get_bibles.get_bible_verses(index_text, chapter, 26, 30)]
            
    ]

    add_line = 8  # 每页最多显示7行经文，超过则添加新行
    
    count = 0
    for text in texts:
        count += 1
        i = text[2]
        bibles = text[4]
        
        '''
        replacements = {
            1: {0: {1: text[0], 2: f" {text[1]}: {text[2]}-{text[3]}"}},
            2: {
                0: {0: str(i), 1: bibles[0] if len(bibles) > 0 else ""},
                1: {0: str(i + 1), 1: bibles[1] if len(bibles) > 1 else ""},
                2: {0: str(i + 2), 1: bibles[2] if len(bibles) > 2 else ""},
                3: {0: str(i + 3) if len(bibles) > 3 else "", 1: bibles[3] if len(bibles) > 3 else ""},
                4: {0: str(i + 4) if len(bibles) > 4 else "", 1: bibles[4] if len(bibles) > 4 else ""},
                5: {0: "", 1: ""}
            }
        }
        '''

        replacements = {
            1: {0: {0: f"{text[0]}", 1: f"{text[1]}章", 2: f"{text[2]}", 3: f"-{text[3]}"}},
            2: {
                0: {0: str(i) , 2: bibles[0] if len(bibles) > 0 else ""},
                1: {0: str(i + 1) if len(bibles) > 1 else "", 1: "", 2: bibles[1] if len(bibles) > 1 else ""},
                2: {0: str(i + 2) if len(bibles) > 2 else "", 1: bibles[2] if len(bibles) > 2 else ""},
                3: {0: str(i + 3) if len(bibles) > 3 else "", 1: bibles[3] if len(bibles) > 3 else ""},
                4: {0: str(i + 4) if len(bibles) > 4 else "", 1: bibles[4] if len(bibles) > 4 else ""},
                5: {0: "", 1: ""},
                6: {0: "", 1: ""},
            }
        }
        '''
        if text[3] - text[2] > add_line:
            for j in range(add_line, text[3] - text[2] + 1):
                replacements[2][add_line][2] += f" \n{str(i + j)}    " + (bibles[j] if len(bibles) > j else "")
        
        set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, replacements, resize = 36)
        if count < len(texts):
            duplicate_slide(output_file, output_file, page_to_modify)
        #page_to_modify += 1  
        '''
         
    #delete_slides(output_file, output_file, [13])  # 删除多余的经文页，保留第一页经文页
