from pptx import Presentation
import os
import get_bibles
import copy
from pptx.util import Pt

def read_pptx(pptx_file):
    """
    读取现有的PPTX文件并返回所有内容信息
    
    Args:
        pptx_file: PPTX文件路径
    
    Returns:
        dict: 包含PPT所有信息的字典
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return None
    
    prs = Presentation(pptx_file)
    ppt_info = {
        'slide_count': len(prs.slides),
        'slides': []
    }
    
    # 遍历所有幻灯片
    for slide_num, slide in enumerate(prs.slides, 1):
        slide_info = {
            'slide_number': slide_num,
            'title': '',
            'shapes': []
        }
        
        # 获取标题
        if slide.shapes.title:
            slide_info['title'] = slide.shapes.title.text
        
        # 遍历所有形状
        for shape_num, shape in enumerate(slide.shapes):
            shape_info = {
                'shape_number': shape_num,
                'type': str(shape.shape_type),
                'has_text': shape.has_text_frame,
                'text': ''
            }
            
            # 获取文本内容
            if shape.has_text_frame:
                text_parts = []
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_parts.append(run.text)
                shape_info['text'] = ''.join(text_parts)
            
            slide_info['shapes'].append(shape_info)
        
        ppt_info['slides'].append(slide_info)
    
    return ppt_info


def update_pptx_text(pptx_file, output_file, replacements):
    """
    修改PPTX文件中的文字
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        replacements: 字典，格式 {'旧文字': '新文字'}
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 遍历所有幻灯片
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 进行替换
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
    
    # 保存
    prs.save(output_file)
    print(f"PPT文件已保存: {output_file}")
    return True


def print_pptx_info(ppt_info):
    """
    打印PPT信息
    
    Args:
        ppt_info: read_pptx函数返回的信息字典
    """
    if not ppt_info:
        return
    
    print(f"总共 {ppt_info['slide_count']} 页")
    print("=" * 60)
    
    for slide_info in ppt_info['slides']:
        print(f"\n第 {slide_info['slide_number']} 页")
        print(f"标题: {slide_info['title']}")
        print(f"形状数量: {len(slide_info['shapes'])}")       
        for shape_info in slide_info['shapes']:
            
            if shape_info['has_text'] and shape_info['text']:
                print(f"  - 文本: {shape_info['text']}...")


def print_pptx_page(ppt_info, page_number):
    """
    打印PPT信息
    
    Args:
        ppt_info: read_pptx函数返回的信息字典
    """
    if not ppt_info:
        return
    
    print(f"总共 {ppt_info['slide_count']} 页")
    print("=" * 60)
    
    for slide_info in ppt_info['slides']:
        if slide_info['slide_number'] != page_number:
            continue
        print(f"\n第 {slide_info['slide_number']} 页")
        print(f"标题: {slide_info['title']}")
        print(f"形状数量: {len(slide_info['shapes'])}")       
        for shape_info in slide_info['shapes']:
            
            if shape_info['has_text'] and shape_info['text']:
                print(f"  - 文本: {shape_info['text']}...")


def update_slide_text(pptx_file, output_file, slide_number, replacements):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        replacements: 字典，格式 {'旧文字': '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    # 遍历该页的所有形状
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # 进行替换
                    for old_text, new_text in replacements.items():
                        #print(f"Before {run.text} → {run.text.replace(old_text, new_text)}")
                        if old_text in run.text:
                            run.text = run.text.replace(old_text, new_text)
                        #print(f"After {run.text}")
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True


def update_multiple_slides(pptx_file, output_file, slide_replacements):
    """
    批量修改多页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_replacements: 字典，格式 {页码: {'旧文字': '新文字'}}
        
    Example:
        slide_replacements = {
            1: {'标题': '新标题'},
            2: {'内容': '新内容'},
            3: {'2025': '2026'}
        }
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 遍历需要修改的页
    for slide_number, replacements in slide_replacements.items():
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 获取指定页
        slide = prs.slides[slide_number - 1]
        
        # 遍历该页的所有形状
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # 进行替换
                        for old_text, new_text in replacements.items():
                            if old_text in run.text:
                                run.text = run.text.replace(old_text, new_text)
        
        print(f"已修改第 {slide_number} 页")
    
    # 保存
    prs.save(output_file)
    print(f"所有修改完成，文件已保存: {output_file}")
    return True


def delete_slide(pptx_file, output_file, slide_number):
    """
    删除指定页
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 要删除的页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取要删除的幻灯片
    rId = prs.slides._sldIdLst[slide_number - 1].rId
    prs.part.drop_rel(rId)
    del prs.slides._sldIdLst[slide_number - 1]
    
    # 保存
    prs.save(output_file)
    print(f"已删除第 {slide_number} 页，文件已保存: {output_file}")
    return True


def delete_slides(pptx_file, output_file, slide_numbers):
    """
    批量删除多页
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_numbers: 要删除的页码列表（从1开始），如 [2, 5, 7]
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 从大到小排序，从后往前删除，避免索引变化
    slide_numbers_sorted = sorted(slide_numbers, reverse=True)
    
    for slide_number in slide_numbers_sorted:
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 删除幻灯片
        rId = prs.slides._sldIdLst[slide_number - 1].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[slide_number - 1]
        print(f"已删除第 {slide_number} 页")
    
    # 保存
    prs.save(output_file)
    print(f"所有删除完成，文件已保存: {output_file}")
    return True


def duplicate_slides(pptx_file, output_file, slide_numbers):
    """
    批量复制多页并插入到各自后面
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_numbers: 要复制的页码列表（从1开始），如 [2, 5]
    
    Returns:
        bool: 是否成功
    """
    import copy
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 从大到小排序，从后往前处理，避免索引变化
    slide_numbers_sorted = sorted(slide_numbers, reverse=True)
    
    for slide_number in slide_numbers_sorted:
        # 检查页码是否有效
        if slide_number < 1 or slide_number > len(prs.slides):
            print(f"警告：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页），跳过")
            continue
        
        # 获取要复制的幻灯片
        source_slide = prs.slides[slide_number - 1]
        
        # 获取布局
        slide_layout = source_slide.slide_layout
        
        # 创建新幻灯片
        new_slide = prs.slides.add_slide(slide_layout)
        
        # 深度复制所有形状
        for shape in source_slide.shapes:
            el = shape.element
            newel = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
        
        # 移动到正确位置
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(slide_number, slides[-1])
        
        print(f"已在第 {slide_number} 页后插入副本")
    
    # 保存
    prs.save(output_file)
    print(f"所有复制完成，文件已保存: {output_file}")
    return True


def show_structure_one_page(pptx_file, slide_number):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        indexed_replacements: 字典，格式 {索引: '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]  
    
    for i in range(len(slide.shapes)):
        shape = slide.shapes[i]
        print(f"Shape index: {i}")
        if slide.shapes[i].has_text_frame:
            for j in range(len(slide.shapes[i].text_frame.paragraphs)):
                paragraph = slide.shapes[i].text_frame.paragraphs[j]
                print(f"  Paragraph index: {j}")
                for k in range(len(paragraph.runs)):
                    run = paragraph.runs[k]
                    print(f"     text index : {k} : {run.text}", end="|\n")


def duplicate_slide(pptx_file, output_file, slide_number):
    """
    复制指定页并插入到该页后面
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 要复制的页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    import copy
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取要复制的幻灯片
    source_slide = prs.slides[slide_number - 1]
    
    # 获取布局
    slide_layout = source_slide.slide_layout
    
    # 创建新幻灯片
    new_slide = prs.slides.add_slide(slide_layout)
    
    # 深度复制所有形状
    for shape in source_slide.shapes:
        el = shape.element
        newel = copy.deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newel, 'p:extLst')
    
    # 移动到正确位置（紧跟在原页面后）
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[-1])
    xml_slides.insert(slide_number, slides[-1])
    
    # 保存
    prs.save(output_file)
    print(f"已在第 {slide_number} 页后插入副本，文件已保存: {output_file}")
    return True


def swap_slides(pptx_file, output_file, slide_num1, slide_num2):
    """
    交换两个幻灯片的位置
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_num1: 第一个页码（从1开始）
        slide_num2: 第二个页码（从1开始）
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_num1 < 1 or slide_num1 > len(prs.slides):
        print(f"错误：页码 {slide_num1} 超出范围（共 {len(prs.slides)} 页）")
        return False
    if slide_num2 < 1 or slide_num2 > len(prs.slides):
        print(f"错误：页码 {slide_num2} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    if slide_num1 == slide_num2:
        print("错误：两个页码不能相同")
        return False
    
    # 转换为0开始的索引
    idx1 = slide_num1 - 1
    idx2 = slide_num2 - 1
    
    # 获取XML中的slides列表
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    
    # 交换位置
    slides[idx1], slides[idx2] = slides[idx2], slides[idx1]
    
    # 清空并重新添加
    for slide in list(xml_slides):
        xml_slides.remove(slide)
    
    for slide in slides:
        xml_slides.append(slide)
    
    # 保存
    prs.save(output_file)
    print(f"已交换第 {slide_num1} 页和第 {slide_num2} 页，文件已保存: {output_file}")
    return True


def insert_fullscreen_video_slide(pptx_file, output_file, video_path, insert_position=None):
    """
    插入一个新的全屏视频幻灯片
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        video_path: 视频文件路径
        insert_position: 插入位置（从1开始），如果为None则在末尾添加
    
    Returns:
        bool: 是否成功
    """
    from pptx.util import Inches
    
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    if not os.path.exists(video_path):
        print(f"错误：找不到视频文件 {video_path}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 获取幻灯片尺寸
    slide_width = prs.slide_width
    slide_height = prs.slide_height
    
    # 添加一个空白幻灯片（使用空白布局）
    blank_slide_layout = prs.slide_layouts[6]  # 6通常是空白布局
    new_slide = prs.slides.add_slide(blank_slide_layout)
    
    # 添加全屏视频
    # 视频位置：左上角(0,0)，尺寸：填满整个幻灯片
    left = Inches(0)
    top = Inches(0)
    width = slide_width
    height = slide_height
    
    # 插入视频
    movie = new_slide.shapes.add_movie(
        video_path,
        left, top, width, height,
        poster_frame_image=None,  # 不使用海报帧，使用视频第一帧
        mime_type='video/mp4'
    )
    
    # 如果指定了插入位置，则移动到该位置
    if insert_position is not None:
        if insert_position < 1 or insert_position > len(prs.slides):
            print(f"错误：插入位置 {insert_position} 超出范围（共 {len(prs.slides)} 页）")
            return False
        
        # 获取XML中的slides列表
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        
        # 移动新添加的幻灯片（最后一个）到指定位置
        xml_slides.remove(slides[-1])
        xml_slides.insert(insert_position - 1, slides[-1])
    
    # 保存
    prs.save(output_file)
    position_str = f"第 {insert_position} 页" if insert_position else "末尾"
    print(f"已在 {position_str} 插入全屏视频幻灯片，文件已保存: {output_file}")
    return True


def set_pptx_page_texts(pptx_file, output_file, slide_number, replacements):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        replacements: 字典，格式 {'旧文字': '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    # 遍历该页的所有形状
    for shape in slide.shapes:
        if shape.has_text_frame:
            for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for origin_text, change_text in replacements.items():
                            if origin_text in run.text:
                                print(f"{run.text} → {change_text}", end="\n")
                                run.text = run.text.replace(origin_text, change_text)
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True


def set_pptx_page_texts_by_slides_shapes_index(pptx_file, output_file, slide_number, replacements):
    """
    修改指定页的文字内容
    
    Args:
        pptx_file: 原PPTX文件路径
        output_file: 输出PPTX文件路径
        slide_number: 页码（从1开始）
        indexed_replacements: 字典，格式 {索引: '新文字'}
    
    Returns:
        bool: 是否成功
    """
    if not os.path.exists(pptx_file):
        print(f"错误：找不到文件 {pptx_file}")
        return False
    
    prs = Presentation(pptx_file)
    
    # 检查页码是否有效
    if slide_number < 1 or slide_number > len(prs.slides):
        print(f"错误：页码 {slide_number} 超出范围（共 {len(prs.slides)} 页）")
        return False
    
    # 获取指定页（索引从0开始）
    slide = prs.slides[slide_number - 1]
    
    for shape_index, run_replacements in replacements.items():
        if not slide.shapes.__getitem__(shape_index).has_text_frame:
            print(f"错误：形状索引 {shape_index} 不包含文本框")
            return False
        shape = slide.shapes[shape_index]
        for paragraph_index in run_replacements.keys():
            paragraph = shape.text_frame.paragraphs[paragraph_index] 
            new_texts_index = run_replacements[paragraph_index]
            print(f"Shape {shape_index} paragraph {paragraph_index}")
            for run_index, new_text in new_texts_index.items():
                if run_index < len(paragraph.runs):
                    print(f" original text {paragraph.runs[run_index].text} new text: {new_text}")
                    paragraph.runs[run_index].text = new_text
                    paragraph.runs[run_index].font.bold = True
                    paragraph.runs[run_index].font.size = Pt(38)
                else:
                    print(f" append new text on {run_index}: {new_text}")
                    new_run = paragraph.add_run()
                    new_run.text = " " + new_text
                    # 新增行时，字体加粗、字号30pt，字体固定为STXingkai
                    new_run.font.bold = True
                    new_run.font.size = Pt(20)
                    new_run.font.name = "STXingkai"
    
    
    # 保存
    prs.save(output_file)
    print(f"已修改第 {slide_number} 页，文件已保存: {output_file}")
    return True

if __name__ == "__main__":
    # 示例1：读取PPT信息
    filename = "template_français"

    template_repo = f'D:\\副业赚钱\\教会事务\\Template\\{filename}.pptx'
    repository = os.path.dirname(os.path.abspath(__file__))
    repository_music = os.path.join(os.path.dirname(repository), 'Template', 'musics')
    print(f"当前路径: {repository}")

    pptx_file = f"{repository}\\{filename}.pptx"
    output_file = f"{repository}\\{filename}.pptx"

    info = read_pptx(output_file)

    #delete_slides(output_file, output_file, [2,3,4])
    
    # 1 时间
    page_to_modify = 1
    date = "01/02/2026 \n"
    heure = "              13h30-14h30\n"
    remplacements = {0: {3: {1: "法语课 Bienvenue !\n"}, 4: {1: date, 2: heure}}, 1: {0: {0: "", 1: "", 2: "", 3: "", 4: "", 5: "", 6: "",  7: "", 8: ""}}}
    #set_pptx_page_texts(output_file, output_file, page_to_modify, {old_date: date}) 
    #set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, remplacements)

    # 2 经文
    page_to_modify = 2
    #duplicate_slide(output_file, output_file, page_to_modify - 1)

    book_zh = "路加福音"      # 路加福音
    book_fr = ""
    chapter_num = 9    # 第1章
    start = 48          # 第1节
    end = 48           # 到第5节

    text = "\n"+ get_bibles.get_bible_verses(book_zh, chapter_num, start, end)[0] 
    text_fr = "\nQuiconque reçoit en mon nom ce petit enfant me reçoit moi-même; et quiconque me reçoit reçoit celui qui m'a envoyé. Car celui qui est le plus petit parmi vous tous, c'est celui-là qui est grand."
    remplacements = {0: {3: {1: text}, 4: {1: "", 2: ""}}}
    #show_structure_one_page(output_file, page_to_modify)
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, remplacements)

    duplicate_slide(output_file, output_file, page_to_modify)
    page_to_modify += 1
    remplacements = {0: {3: {1: text_fr}, 4: {1: "", 2: ""}}}
    set_pptx_page_texts_by_slides_shapes_index(output_file, output_file, page_to_modify, remplacements)

    # 3 诗歌
    music = f"{repository_music}\\4.mp4"
    #insert_fullscreen_video_slide(output_file, output_file, music, insert_position=3)


    